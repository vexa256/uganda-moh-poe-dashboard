#!/usr/bin/env python3
"""
build-tot-pptx.py — generate the Uganda POE Train-the-Trainer deck.

Premium 16:9 deck. Cover, two-dashboard intro, reading-the-page anatomy,
22 surface walkthroughs (one per slide with embedded screenshot + bullet
points + full speaker notes), section dividers, demo-discipline rules,
closing slide.

Run with:
    python3 docs/training/scripts/build-tot-pptx.py
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text  import PP_ALIGN, MSO_ANCHOR
from PIL import Image
Image.MAX_IMAGE_PIXELS = None  # accept full-page screenshots at 2× DPI

# ─── Paths ────────────────────────────────────────────────────────────────
ROOT     = Path('/home/hacker/ecsa-uganda-poe')
SHOTS    = ROOT / 'docs' / 'training' / 'screenshots'
OUT_PPTX = ROOT / 'docs' / 'training' / 'Uganda-POE-Dashboard-ToT.pptx'

# ─── Design tokens ────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0B, 0x25, 0x45)   # primary
NAVY_DARK = RGBColor(0x06, 0x14, 0x28)
TEAL      = RGBColor(0x00, 0xB4, 0xA6)   # accent
TEAL_DARK = RGBColor(0x00, 0x80, 0x76)
RED       = RGBColor(0xD8, 0x31, 0x5B)
AMBER     = RGBColor(0xF4, 0xA8, 0x26)
INK       = RGBColor(0x11, 0x18, 0x27)
SLATE     = RGBColor(0x4B, 0x55, 0x63)
PAPER     = RGBColor(0xF7, 0xFA, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RULE      = RGBColor(0xE5, 0xE7, 0xEB)

FONT      = 'Calibri'

# ─── Surface dictionary — one entry per slide ─────────────────────────────
# (slug, title, subtitle, bullets, speaker_notes)
QR_BULLETS_ANATOMY = [
    'Filters at the top — date range, POE, status, search.',
    'KPI cards — five to eight headline numbers.',
    'One adaptive bar chart — it picks the cleanest view for the data.',
    '20-row table — full cohort in the CSV export.',
]

SURFACES = [
    # ── 0 · landing ──────────────────────────────────────────────────────
    {
        'slug':  '00-screening-volume',
        'title': 'Screening Volume',
        'kicker': 'Quick Report · /admin/quick-reports/screening-volume',
        'family': 'Quick Reports',
        'bullets': [
            'Total primary vs secondary screenings in the date window.',
            'Gender split (M / F / Other / Unknown) and median age band.',
            'Escalation % — the fraction of primary screenings that became referrals.',
            'This is the post-login landing — every user arrives here.',
        ],
        'notes':
            'Open with this every session — it is what every user sees on first sign-in. '
            'Read the first KPI aloud (total primary), then the escalation %. '
            'Point at the chart and read its axis labels. '
            'Skip the date range "all time" — keep filters at Past 7 days while teaching.',
    },
    {
        'slug':  '01-suspected-cases',
        'title': 'Suspected Cases',
        'kicker': 'Quick Report · /admin/quick-reports/suspected-cases',
        'family': 'Quick Reports',
        'bullets': [
            'The clinical priority report — open here when asked "who should we be looking at right now?"',
            'Five KPIs: total, with disease, high-risk, open cases, last 24 hours.',
            'Table rows are deep links to the case file at /admin/alerts/{id}/case-file.',
            'The placeholder "no_specific_suspicion" is quietly dropped from charts so counts stay clinically honest.',
        ],
        'notes':
            'This is the most-asked-for report. Teach it second only to Screening Volume. '
            'Click one row to drop into the case file — that is the deepest view in the system. '
            'Back out to return to the report. The chart reshapes itself; do not promise a fixed shape.',
    },
    {
        'slug':  '02-confirmed-cases',
        'title': 'Confirmed Cases',
        'kicker': 'Quick Report · /admin/quick-reports/confirmed-cases',
        'family': 'Quick Reports',
        'bullets': [
            'The lab pipeline view — confirmed vs probable vs ruled-out vs pending.',
            'Six KPIs: total, confirmed, probable, ruled out, pending, last 24 hours.',
            'Read this alongside Suspected Cases — suspected is the engine guess, confirmed is the lab truth.',
            'Filter by IHR tier to see only internationally-notifiable confirmations.',
        ],
        'notes':
            'Always teach this immediately after Suspected Cases — the pair tells the full clinical story. '
            'Read one KPI from each side aloud (suspected open vs confirmed) so the contrast is felt.',
    },
    {
        'slug':  '03-alert-database',
        'title': 'Alert Database',
        'kicker': 'Quick Report · /admin/quick-reports/alert-database',
        'family': 'Quick Reports',
        'bullets': [
            'The full alert ledger — every alert ever, open or closed.',
            'Seven KPIs: total, open, acknowledged, closed, reopened, tier-1, last 24 hours.',
            'Owner column shows who is handling each alert.',
            'Sorting and filters are persistent in the URL — bookmark a filtered view to come back to it.',
        ],
        'notes':
            'Use this when an auditor or PHEOC officer asks for the full alert history. '
            'Do NOT switch the date range to "All time" on the demo — on the production-sized server the table visibly loads. '
            'Keep it at Past 30 days.',
    },
    {
        'slug':  '04-alert-analysis',
        'title': 'Alert Analysis',
        'kicker': 'Quick Report · /admin/quick-reports/alert-analysis',
        'family': 'Quick Reports',
        'bullets': [
            'Spotlight on HIGH / CRITICAL risk and IHR tier-1 alerts.',
            'Six KPIs include false-positive rate and last-24h count.',
            'The table is pre-filtered to high-stakes rows; the CSV export carries the full cohort.',
            'Open this for the weekly PHEOC briefing — it surfaces what needs eyes.',
        ],
        'notes':
            'This is the report a supervisor opens to find what their team must follow up on this week. '
            'The chart often shows daily risk distribution — read the dominant colour of each day to read the week\'s mood.',
    },
    {
        'slug':  '05-alert-outcomes',
        'title': 'Alert Outcomes',
        'kicker': 'Quick Report · /admin/quick-reports/alert-outcomes',
        'family': 'Quick Reports',
        'bullets': [
            'SLA performance — median acknowledge time, median closure time.',
            'Closed-within-24h % and closed-within-7d % at the top.',
            'Reopened count and un-acknowledged count flagged in red where relevant.',
            'Use this when leadership asks "how fast are we responding?"',
        ],
        'notes':
            'This is the report that goes into the monthly review deck. '
            'Read the median ack-time in minutes and median close-time in hours aloud — '
            'those two numbers are the team\'s heartbeat.',
    },
    {
        'slug':  '06-symptom-spread',
        'title': 'Symptom Spread',
        'kicker': 'Quick Report · /admin/quick-reports/symptom-spread',
        'family': 'Quick Reports',
        'bullets': [
            'Clinical symptom surveillance — what symptoms are appearing and where.',
            'Red-flag symptoms highlighted in red in the chart and in the table.',
            'Top symptom KPI surfaces the most common signal for the date window.',
            'Use this for the morning surveillance briefing.',
        ],
        'notes':
            'Teach this as the "what is going around" report. The red-flag highlighting is automatic — '
            'no settings to configure. Read the top symptom and the red-flag case count aloud.',
    },
    {
        'slug':  '07-poe-analysis',
        'title': 'POE Analysis',
        'kicker': 'Quick Report · /admin/quick-reports/poe-analysis',
        'family': 'Quick Reports',
        'bullets': [
            'Workload across every border — which POEs are busy, which are dark.',
            'Active POEs count, dark POEs count, alert rate as a percentage.',
            'Last-activity column shows when each POE last reported.',
            'Open this when comparing borders side by side.',
        ],
        'notes':
            'This is the report that decides which border needs attention this week. '
            'A "dark POE" is one with no recent activity — that is information, not necessarily a problem.',
    },
    {
        'slug':  '08-country-analysis',
        'title': 'Country Analysis',
        'kicker': 'Quick Report · /admin/quick-reports/country-analysis',
        'family': 'Quick Reports',
        'bullets': [
            'Travel epidemiology — nationalities, visited countries, transit routes.',
            'Endemic-country flow highlighted in red where it matches the disease list.',
            'Seven KPIs include distinct visited countries and endemic case count.',
            'Use this for any travel-history question.',
        ],
        'notes':
            'Teach the three lenses: where they are from, where they have been, where they transited. '
            'The red endemic highlight is automatic — derived from the disease-to-country mapping.',
    },
    {
        'slug':  '09-daily-screening',
        'title': 'Daily Screening',
        'kicker': 'Quick Report · /admin/quick-reports/daily-screening',
        'family': 'Quick Reports',
        'bullets': [
            'The daily throughput view — volume trend, escalation rate, gender split.',
            'Today vs yesterday vs 7-day average KPIs at the top.',
            'Busiest day and busiest POE KPIs surface the operational peak.',
            'Open this for daily ops standup.',
        ],
        'notes':
            'This is the report that opens the daily standup. '
            'Read today\'s primary screenings, then escalation %, then the busiest POE.',
    },
    {
        'slug':  '10-user-analysis',
        'title': 'User Analysis',
        'kicker': 'Quick Report · /admin/quick-reports/user-analysis',
        'family': 'Quick Reports',
        'bullets': [
            'Officer activity — active, dormant, inactive, with workload counts each.',
            'Top performer KPI surfaces the leader for the date window.',
            'Median screenings per active officer = the per-person workload metric.',
            'Use this for monthly workforce reviews.',
        ],
        'notes':
            'Teach this once a month. The chart often ranks officers by screenings; '
            'the leader is shown in red so the eye finds them immediately.',
    },

    # ── Alert Lifecycle ──────────────────────────────────────────────────
    {
        'slug':  '11-alerts-hub',
        'title': 'Alerts — the workflow hub',
        'kicker': 'Alert Lifecycle · /admin/alerts',
        'family': 'Alert Lifecycle',
        'bullets': [
            'The single workflow surface for every alert — Open, Acknowledged, Closed, Reopened, All.',
            'Five tabs at the top; "New" (Open) is the default landing.',
            'Click the 📋 case-file icon (right-side action strip) to open the deep case file.',
            'Click anywhere else on the row to open the gateway modal — close it without acting during a demo.',
        ],
        'notes':
            'Demo discipline: never click an Action button (Acknowledge / Close / Reassign) during teaching. '
            'Walk to the gateway modal, name what you see, close it. Open the case file via the 📋 icon only — '
            'the case file is read-only, safe to demo.',
    },

    # ── Workforce ────────────────────────────────────────────────────────
    {
        'slug':  '12-workforce',
        'title': 'Workforce',
        'kicker': 'Workforce · /admin/workforce',
        'family': 'Workforce',
        'bullets': [
            'One unified page for Users, Roles, and Assignments — three tabs.',
            '+ Add person button opens the atomic four-step wizard.',
            'Only NATIONAL_ADMIN can create users; other roles read their scope.',
            'Per-row icons: edit, suspend, reset password, regenerate invite, disable.',
        ],
        'notes':
            'Walk the wizard four steps without submitting: Identity → Role → Jurisdiction → Invite Mode. '
            'Pick "Credential" not "Email" during demos so nothing leaves the room. '
            'Test users should be named "Training Test 1 — [initials]" so they are easy to disable afterwards.',
    },

    # ── Geography ────────────────────────────────────────────────────────
    {
        'slug':  '13-geo-regions',
        'title': 'Regions',
        'kicker': 'Geography · /admin/geo/provinces',
        'family': 'Geography',
        'bullets': [
            'The regional PHEOCs — the country\'s coordination units.',
            'Each region shows District count and POE count — what sits underneath.',
            'NATIONAL_ADMIN writes; everyone else reads.',
            'Renaming a region cascades to every POE\'s region field and bumps the bundle version.',
        ],
        'notes':
            'Teach the cascading rename rule by saying it, not by doing it. '
            'A rename triggers every mobile device to re-sync — recoverable on training, painful on production.',
    },
    {
        'slug':  '14-geo-districts',
        'title': 'Districts',
        'kicker': 'Geography · /admin/geo/districts',
        'family': 'Geography',
        'bullets': [
            'Border-adjacent districts. Each district sits under one region.',
            'Name "_raw" auto-strips the " District" suffix for the mobile payload.',
            'Cannot delete a district that still has POEs or hospitals under it.',
            'Filter by region in the filter bar to focus the table.',
        ],
        'notes':
            'Same shape as Regions. The cannot-delete-with-children rule is the one new teaching point — '
            'state it once and move on.',
    },
    {
        'slug':  '15-geo-hospitals',
        'title': 'Hospitals',
        'kicker': 'Geography · /admin/geo/hospitals',
        'family': 'Geography',
        'bullets': [
            'The referral hospital roster — type (teaching / general / district / rural / clinic / private / military).',
            'Each hospital sits under one region and one district.',
            'Greenfield surface today — the table is intentionally light.',
            'No cascading, no bundle version bump on writes here.',
        ],
        'notes':
            'Mention this is a greenfield surface — the dataset is still being populated. '
            'Do not over-explain a sparse page; show the columns and move on.',
    },
    {
        'slug':  '16-geo-countries',
        'title': 'Countries',
        'kicker': 'Geography · /admin/geo/countries',
        'family': 'Geography',
        'bullets': [
            'Single-tenant — one row, Uganda.',
            'Edit the display name, ISO codes, and metadata without a code change.',
            'Province / district / POE counts shown on the row.',
            'Writes bump the bundle version.',
        ],
        'notes':
            'A one-row surface — show that the row exists, that it has counts of everything under it, '
            'then move on. Do not edit during the demo.',
    },

    # ── PoEs · Annex-1A ──────────────────────────────────────────────────
    {
        'slug':  '17-poe-registry',
        'title': 'PoE Registry',
        'kicker': 'PoEs · Annex-1A · /admin/geo/poes',
        'family': 'PoEs · Annex-1A',
        'bullets': [
            'Every gazetted Point of Entry — airports, ports, land borders, island entries, rail crossings.',
            'Create / edit is a five-step wizard — Identity, Location, Flags, Context, Review.',
            'Smart auto-fill suggests POE type, transport mode and a unique external ID as you type.',
            'Duplicate-check pre-flights before save and shows a "this looks like an existing POE" warning.',
        ],
        'notes':
            'Walk the five wizard steps without saving — every save bumps the bundle version and every phone re-syncs. '
            'Show that typing a POE name triggers the suggestion panel. Stop at Step 5 — Review.',
    },
    {
        'slug':  '18-poe-capacity',
        'title': 'PoE Capacity',
        'kicker': 'PoEs · Annex-1A · /admin/poe/capacity',
        'family': 'PoEs · Annex-1A',
        'bullets': [
            'WHO IHR Annex-1A capacity scoring — eight dimensions, each scored 1 to 5.',
            'One-way workflow: DRAFT → SUBMITTED → REVIEWED. Once submitted, no edits.',
            'Three tabs at the top — DRAFT, SUBMITTED, REVIEWED.',
            'Overall score is computed from the eight dimensions: ((avg − 1) ÷ 4) × 100.',
        ],
        'notes':
            'Show all three tabs so trainees see the one-way ratchet. Click into one assessment to expose the '
            'eight-dimension scoring panel. Do not move an assessment forward during the demo.',
    },
    {
        'slug':  '19-poe-status',
        'title': 'PoE Status',
        'kicker': 'PoEs · Annex-1A · /admin/poe/status',
        'family': 'PoEs · Annex-1A',
        'bullets': [
            'Time-series log of status changes — Open, Closed, Reduced Hours, Emergency Closed, Maintenance.',
            'The latest row with no end-date is the "current" status for each POE.',
            'Posting a new status automatically closes the previous one — you record next, you do not edit current.',
            'Recent log shows up to 30 events below the current row.',
        ],
        'notes':
            'Teach the "record next" mental model — this is not a status field, it is a status log. '
            'Show the current row for one POE, then the history below it. Do not post during the demo.',
    },
    {
        'slug':  '20-roster-ladder',
        'title': 'Roster & Ladder',
        'kicker': 'PoEs · Annex-1A · /admin/poe/contacts',
        'family': 'PoEs · Annex-1A',
        'bullets': [
            'Contact list per POE — name, organisation, phone, email, channel.',
            'Ten receives_* flags — Critical, High, Medium, Low, Tier-1, Tier-2, Breach, Followup, Daily, Weekly.',
            'Five-level escalation ladder — POE → District → PHEOC → National → WHO.',
            'Cycle detection — a contact cannot escalate to themselves or to anyone in their own chain.',
        ],
        'notes':
            'Walk through one contact row top-to-bottom: name, level, organisation, phone, email, the receives flags, '
            'who they escalate to. The ladder is the most important concept on this page — '
            'point at the Escalates-To column and read the chain aloud.',
    },
]

# ─── Helpers ──────────────────────────────────────────────────────────────
def add_filled_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, bullets, *, size=14, color=INK, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = '•  ' + line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def set_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(11)

def fit_image(path, target_w_in, target_h_in):
    """Return (w_in, h_in) preserving aspect, fitted into the target box."""
    with Image.open(path) as im:
        iw, ih = im.size
    ar_img = iw / ih
    ar_box = target_w_in / target_h_in
    if ar_img > ar_box:
        return target_w_in, target_w_in / ar_img
    return target_h_in * ar_img, target_h_in

def crop_top_to_aspect(src_path, dst_path, target_aspect):
    """Crop the screenshot to its top portion at target_aspect (w/h).

    Full-page screenshots are very tall (the page scrolls); cropping to the
    top viewport preserves what matters (header → KPIs → chart top) and
    keeps the embedded image at a slide-friendly aspect ratio.
    """
    with Image.open(src_path) as im:
        iw, ih = im.size
        target_h = int(iw / target_aspect)
        if target_h >= ih:
            # Image is already shorter than the target — pad if needed, else
            # just use as is.
            return src_path
        cropped = im.crop((0, 0, iw, target_h))
        # Downsize to a sensible width to keep PPTX small.
        max_w = 2200
        if cropped.size[0] > max_w:
            scale = max_w / cropped.size[0]
            cropped = cropped.resize(
                (max_w, int(cropped.size[1] * scale)),
                Image.LANCZOS,
            )
        cropped.save(dst_path, 'PNG', optimize=True)
    return dst_path

# ─── Build ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height

# ── COVER ────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_filled_rect(s, 0, 0, SW, SH, NAVY_DARK)
# diagonal teal accent
add_filled_rect(s, Inches(0), Inches(6.1), SW, Inches(0.18), TEAL)
# Pre-title chip
add_text(s, Inches(0.85), Inches(0.85), Inches(4), Inches(0.4),
         'UGANDA · POINTS OF ENTRY · 2026', size=12, bold=True, color=TEAL)
# Title
add_text(s, Inches(0.85), Inches(1.3), Inches(11.5), Inches(2.0),
         'Web Dashboard', size=58, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(2.55), Inches(11.5), Inches(1.4),
         'Train-the-Trainer Briefing', size=44, bold=False, color=WHITE)
# Lede
add_text(s, Inches(0.85), Inches(4.3), Inches(10.5), Inches(1.6),
         'Every surface, taught the same way: what it is · who opens it · what is on the screen · the walk · '
         'the hands-on · the one thing to skip during the demo.',
         size=18, color=PAPER)
# Footer chip
add_text(s, Inches(0.85), Inches(6.5), Inches(8), Inches(0.4),
         'For facilitators only · Training dashboard · 2026-05-20', size=12, color=PAPER)
add_text(s, Inches(10.4), Inches(6.5), Inches(2.5), Inches(0.4),
         'v1.0', size=12, color=PAPER, align=PP_ALIGN.RIGHT)
set_notes(s,
    'Open the session by reading the title and lede aloud. '
    'Set expectations: this deck is a teach-script, not a slide-show. '
    'Every screen is taught the same way, and the trainees do something hands-on on every screen.')

# ── INTRO 1 — Two dashboards ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_filled_rect(s, 0, 0, SW, Inches(1.05), NAVY)
add_text(s, Inches(0.85), Inches(0.32), Inches(12), Inches(0.6),
         'Two dashboards. One rule.', size=28, bold=True, color=WHITE)
# Two columns
COL_W = Inches(5.6)
COL_H = Inches(5.2)
COL_Y = Inches(1.55)
COL_GAP = Inches(0.4)

# Training (left, highlighted)
LX = Inches(0.85)
add_filled_rect(s, LX, COL_Y, COL_W, COL_H, PAPER)
add_filled_rect(s, LX, COL_Y, COL_W, Inches(0.6), TEAL)
add_text(s, LX + Emu(91440), COL_Y + Emu(45720), COL_W - Emu(182880), Inches(0.6),
         'TRAINING — we use this in the room', size=14, bold=True, color=WHITE)
add_text(s, LX + Emu(91440), COL_Y + Inches(0.85), COL_W - Emu(182880), Inches(0.5),
         'ug-poe.ecsahc.com/admin', size=16, bold=True, color=NAVY, font='Consolas')
add_bullets(s, LX + Emu(91440), COL_Y + Inches(1.55), COL_W - Emu(182880), Inches(3.4), [
    'Safe sandbox — practise freely, nothing is real.',
    'Pre-seeded with travellers, alerts, POEs and contacts.',
    'Identical layout to the live dashboard — same data shape.',
    'Reset nightly at 23:00 EAT.',
], size=14, color=SLATE)

# Live (right, muted)
RX = LX + COL_W + COL_GAP
add_filled_rect(s, RX, COL_Y, COL_W, COL_H, PAPER)
add_filled_rect(s, RX, COL_Y, COL_W, Inches(0.6), SLATE)
add_text(s, RX + Emu(91440), COL_Y + Emu(45720), COL_W - Emu(182880), Inches(0.6),
         'LIVE — never in the training room', size=14, bold=True, color=WHITE)
add_text(s, RX + Emu(91440), COL_Y + Inches(0.85), COL_W - Emu(182880), Inches(0.5),
         'poes.health.go.ug/admin', size=16, bold=True, color=NAVY, font='Consolas')
add_bullets(s, RX + Emu(91440), COL_Y + Inches(1.55), COL_W - Emu(182880), Inches(3.4), [
    'Real travellers. Real alerts. Real consequences.',
    'Day-to-day operations only — never used in training.',
    'Same screens, same buttons, same workflow as the training one.',
    'If the URL is not the training URL, close the tab.',
], size=14, color=SLATE)

# Footer rule
add_text(s, Inches(0.85), Inches(6.95), Inches(11.5), Inches(0.4),
         'Everything you teach today is identical on the live dashboard. Only the data is different.',
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
set_notes(s,
    'Read the two URLs aloud — slowly. Make sure trainees can hear the difference. '
    'Press the point that everything is identical between the two — same colours, same buttons, same flow. '
    'Only the data behind it is different. Then commit: today, only the left side. '
    'If anyone sees the right URL on a facilitator screen during the session, raise a hand immediately.')

# ── INTRO 2 — Anatomy of a page ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_filled_rect(s, 0, 0, SW, Inches(1.05), NAVY)
add_text(s, Inches(0.85), Inches(0.32), Inches(12), Inches(0.6),
         'How to read every page', size=28, bold=True, color=WHITE)
# Anatomy diagram (left) + text (right)
DX = Inches(0.85); DY = Inches(1.55); DW = Inches(6.8); DH = Inches(5.4)
# wrapper
add_filled_rect(s, DX, DY, DW, DH, PAPER)
# rows
rowH = Inches(0.95)
gap  = Inches(0.18)
rows = [
    ('Filter bar',  'date range · POE · status · search',           TEAL),
    ('KPI cards',   '5–8 headline numbers, important first',        NAVY),
    ('One chart',   'bar chart that picks the cleanest view',       AMBER),
    ('One table',   '20 rows on screen · full cohort in CSV export',RED),
]
ry = DY + Inches(0.25)
for label, sub, col in rows:
    add_filled_rect(s, DX + Inches(0.25), ry, DW - Inches(0.5), rowH, col)
    add_text(s, DX + Inches(0.45), ry + Inches(0.15), DW - Inches(0.9), Inches(0.35),
             label, size=15, bold=True, color=WHITE)
    add_text(s, DX + Inches(0.45), ry + Inches(0.5), DW - Inches(0.9), Inches(0.4),
             sub, size=12, color=WHITE)
    ry += rowH + gap

# Right column — narrative
RX = DX + DW + Inches(0.4)
RW = SW - RX - Inches(0.85)
add_text(s, RX, DY, RW, Inches(0.5),
         'The four-finger walk', size=22, bold=True, color=NAVY)
add_bullets(s, RX, DY + Inches(0.65), RW, Inches(4.5), [
    'Filters at the top — the question you can change.',
    'KPIs underneath — the headline numbers, left to right.',
    'The chart — the picture of the headline. Read the axis labels aloud.',
    'The table — the names behind the numbers. Click a row to see the case file.',
    '',
    'Teach this shape once. The same shape repeats on eleven reports.',
], size=14, color=SLATE)
set_notes(s,
    'This is the most important slide of the day. Eleven of the twenty-one surfaces have this exact shape. '
    'If a trainee learns to read one Quick Report, they can read any of them. '
    'Spend three minutes here — point at each layer in turn, name what it does, repeat.')

# ── SECTION DIVIDERS + SURFACES ──────────────────────────────────────────
def add_section_divider(title, subtitle):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_filled_rect(s, 0, 0, SW, SH, NAVY)
    add_filled_rect(s, 0, Inches(3.6), SW, Inches(0.04), TEAL)
    add_text(s, Inches(0.85), Inches(2.6), Inches(12), Inches(0.6),
             title.upper(), size=18, bold=True, color=TEAL)
    add_text(s, Inches(0.85), Inches(3.85), Inches(12), Inches(1.6),
             subtitle, size=42, bold=True, color=WHITE)
    set_notes(s, f'Section divider — {title}. Pause for ten seconds before continuing. '
                 'Let trainees re-set attention.')
    return s

def add_surface_slide(meta):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    # Header band
    add_filled_rect(s, 0, 0, SW, Inches(1.05), NAVY)
    # Family chip on header
    add_text(s, Inches(0.85), Inches(0.18), Inches(8), Inches(0.32),
             meta['family'].upper(), size=11, bold=True, color=TEAL)
    add_text(s, Inches(0.85), Inches(0.46), Inches(11), Inches(0.55),
             meta['title'], size=26, bold=True, color=WHITE)
    add_text(s, Inches(10.2), Inches(0.55), Inches(3), Inches(0.4),
             meta['kicker'].split(' · ')[-1], size=10, color=PAPER,
             align=PP_ALIGN.RIGHT, font='Consolas')

    # Body — left column: screenshot. Right column: bullets.
    LEFT_X = Inches(0.4)
    LEFT_Y = Inches(1.3)
    LEFT_W = Inches(8.2)
    LEFT_H = Inches(5.95)
    # frame
    add_filled_rect(s, LEFT_X, LEFT_Y, LEFT_W, LEFT_H, RULE)
    img_path = SHOTS / f"{meta['slug']}.png"
    if img_path.exists():
        # Crop full-page screenshot to a slide-friendly aspect, then embed.
        pad = Emu(45720)
        max_w = LEFT_W - 2 * pad
        max_h = LEFT_H - 2 * pad
        slot_aspect = (max_w / max_h)
        cropped_path = SHOTS / f"_crop_{meta['slug']}.png"
        crop_top_to_aspect(str(img_path), str(cropped_path), slot_aspect)
        use_path = cropped_path if cropped_path.exists() else img_path
        w_in, h_in = fit_image(str(use_path),
                               max_w / 914400, max_h / 914400)
        img_w = Inches(w_in)
        img_h = Inches(h_in)
        ix = LEFT_X + (LEFT_W - img_w) / 2
        iy = LEFT_Y + (LEFT_H - img_h) / 2
        s.shapes.add_picture(str(use_path), ix, iy, img_w, img_h)

    # Right column
    RX = LEFT_X + LEFT_W + Inches(0.25)
    RY = LEFT_Y
    RW = SW - RX - Inches(0.4)
    add_text(s, RX, RY, RW, Inches(0.5),
             'What to teach', size=16, bold=True, color=NAVY)
    add_bullets(s, RX, RY + Inches(0.55), RW, Inches(4.8),
                meta['bullets'], size=12, color=INK, line_spacing=1.25)

    # Skip-during-demo callout at bottom right
    callY = RY + Inches(5.0)
    add_filled_rect(s, RX, callY, RW, Inches(0.95), AMBER)
    add_text(s, RX + Inches(0.15), callY + Inches(0.1), RW - Inches(0.3), Inches(0.3),
             'DEMO DISCIPLINE', size=10, bold=True, color=NAVY)
    skip_text = meta['notes'].split('. ')[-1] if meta['notes'] else 'Read, never write, during a demo.'
    add_text(s, RX + Inches(0.15), callY + Inches(0.4), RW - Inches(0.3), Inches(0.5),
             skip_text, size=11, color=NAVY)

    set_notes(s, meta['notes'])
    return s

# Group surfaces by family, with dividers
def family_for(idx):
    return SURFACES[idx]['family']

last_family = None
for meta in SURFACES:
    if meta['family'] != last_family:
        subtitle_map = {
            'Quick Reports': 'Eleven one-question reports. One shape, eleven questions.',
            'Alert Lifecycle': 'The single workflow surface — Open · Acknowledged · Closed · Reopened.',
            'Workforce': 'One unified page. People, roles, and assignments.',
            'Geography': 'The reference hierarchy. Country · Region · District · Hospital.',
            'PoEs · Annex-1A': 'Borders as IHR health objects. Registry · Capacity · Status · Roster.',
        }
        add_section_divider(meta['family'], subtitle_map.get(meta['family'], ''))
        last_family = meta['family']
    add_surface_slide(meta)

# ── Demo discipline slide ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_filled_rect(s, 0, 0, SW, Inches(1.05), NAVY)
add_text(s, Inches(0.85), Inches(0.32), Inches(12), Inches(0.6),
         'Demo discipline — the five rules', size=28, bold=True, color=WHITE)

rules = [
    ('1', 'Read, never write, during a demo.',
        'Walk to the form, name what you see, close the form. Real writes happen in the hands-on minutes, not the walk minutes.'),
    ('2', 'Walk in the order taught in this deck.',
        'Quick Reports → Alerts → Workforce → Geography → PoEs. The order is the map.'),
    ('3', 'Use the training URL only.',
        'ug-poe.ecsahc.com/admin — every time. If a live URL ever appears on a facilitator screen, close the tab.'),
    ('4', 'The chart adapts; do not promise a fixed chart shape.',
        'If a trainee says it looked different yesterday, say "the chart picks the cleanest version of the question for the data in view."'),
    ('5', 'One thing per click.',
        'Click → name what changed → wait for a hand to go up → next click. Never two clicks in quick succession on stage.'),
]
rule_y = Inches(1.4)
rule_h = Inches(1.05)
gap    = Inches(0.1)
for num, head, body in rules:
    # Number badge
    add_filled_rect(s, Inches(0.85), rule_y, Inches(0.95), rule_h, NAVY)
    add_text(s, Inches(0.85), rule_y, Inches(0.95), rule_h,
             num, size=42, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Body
    add_text(s, Inches(2.0), rule_y + Inches(0.05), Inches(11), Inches(0.45),
             head, size=18, bold=True, color=NAVY)
    add_text(s, Inches(2.0), rule_y + Inches(0.5), Inches(11), Inches(0.55),
             body, size=13, color=SLATE)
    rule_y += rule_h + gap
set_notes(s,
    'Five rules. Read them aloud once at the very start of the session, then every facilitator commits to them. '
    'Anyone who can quote rules 1, 3 and 5 in order has internalised the deck.')

# ── Closing slide ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_filled_rect(s, 0, 0, SW, SH, NAVY_DARK)
add_filled_rect(s, 0, Inches(6.1), SW, Inches(0.18), TEAL)
add_text(s, Inches(0.85), Inches(2.6), Inches(12), Inches(0.5),
         'YOU ARE READY', size=14, bold=True, color=TEAL)
add_text(s, Inches(0.85), Inches(3.1), Inches(12), Inches(1.6),
         'Now go teach the room.', size=52, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(4.6), Inches(12), Inches(1.2),
         'Same dashboard everywhere. Same shape on every page. Same six-part teach pattern, every surface.',
         size=18, color=PAPER)
add_text(s, Inches(0.85), Inches(6.5), Inches(8), Inches(0.4),
         'Questions during the session: raise it on the round-table at 15:45.',
         size=12, color=PAPER)
add_text(s, Inches(10.4), Inches(6.5), Inches(2.5), Inches(0.4),
         '— end —', size=12, color=PAPER, align=PP_ALIGN.RIGHT)
set_notes(s,
    'Closing slide. Read the call-to-action aloud. '
    'Remind facilitators that questions land at 15:45 — protect the schedule between now and then.')

# ── Save ────────────────────────────────────────────────────────────────
OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPTX)
print(f'Wrote {OUT_PPTX}')
print(f'Slides: {len(prs.slides)}')
print(f'Size:   {OUT_PPTX.stat().st_size / 1024:.0f} KB')
