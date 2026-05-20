#!/usr/bin/env python3
"""
build-user-manual.py — generate the Uganda POE Web Dashboard User Manual.

This is the *user-facing* manual (not the ToT teach-script).  Every screen
is walked end-to-end with annotated screenshots, numbered callouts, plain-
language explanations, and a "What to expect" strip.  An end-user can read
it cover-to-cover and operate every part of the dashboard on their own.

Build with:
    python3 docs/training/scripts/build-user-manual.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text  import PP_ALIGN, MSO_ANCHOR
from PIL import Image
Image.MAX_IMAGE_PIXELS = None

# ── Paths ─────────────────────────────────────────────────────────────────
ROOT     = Path('/home/hacker/ecsa-uganda-poe')
SHOTS    = ROOT / 'docs' / 'training' / 'screenshots'
OUT_PPTX = ROOT / 'docs' / 'training' / 'Uganda-POE-Dashboard-User-Manual.pptx'

# ── Design tokens ─────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0B, 0x25, 0x45)
NAVY_DARK = RGBColor(0x06, 0x14, 0x28)
TEAL      = RGBColor(0x00, 0xB4, 0xA6)
TEAL_DK   = RGBColor(0x00, 0x80, 0x76)
RED       = RGBColor(0xD8, 0x31, 0x5B)
AMBER     = RGBColor(0xF4, 0xA8, 0x26)
INK       = RGBColor(0x11, 0x18, 0x27)
SLATE     = RGBColor(0x4B, 0x55, 0x63)
SLATE_LT  = RGBColor(0x9C, 0xA3, 0xAF)
PAPER     = RGBColor(0xF7, 0xFA, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RULE      = RGBColor(0xE5, 0xE7, 0xEB)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
FONT      = 'Calibri'

# ── Low-level helpers ─────────────────────────────────────────────────────
def rect(slide, x, y, w, h, color, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = color
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def round_rect(slide, x, y, w, h, color, *, radius=Pt(8)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.16
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def text(slide, x, y, w, h, txt, *, size=14, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = txt
    f = run.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return tb

def bullets(slide, x, y, w, h, items, *, size=12, color=INK,
            line=1.25, marker='•  '):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line_text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line
        run = p.add_run()
        run.text = (marker if marker else '') + line_text
        f = run.font; f.name = FONT; f.size = Pt(size); f.color.rgb = color
    return tb

def numbered(slide, x, y, w, h, items, *, size=12, num_color=TEAL, text_color=INK):
    """Numbered list — used for callout legends."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, txt_line in enumerate(items, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.30
        run = p.add_run()
        run.text = f'{i}.  '
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = True
        run.font.color.rgb = num_color
        run = p.add_run()
        run.text = txt_line
        run.font.name = FONT; run.font.size = Pt(size); run.font.color.rgb = text_color
    return tb

def notes(slide, txt):
    nt = slide.notes_slide.notes_text_frame
    nt.text = txt
    for p in nt.paragraphs:
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(11)

def crop_top(src, dst, target_aspect, *, max_w=2200):
    with Image.open(src) as im:
        iw, ih = im.size
        target_h = int(iw / target_aspect)
        if target_h < ih:
            im = im.crop((0, 0, iw, target_h))
        if im.size[0] > max_w:
            scale = max_w / im.size[0]
            im = im.resize((max_w, int(im.size[1] * scale)), Image.LANCZOS)
        im.save(dst, 'PNG', optimize=True)

def fit(path, box_w_in, box_h_in):
    with Image.open(path) as im:
        iw, ih = im.size
    ar_i = iw / ih; ar_b = box_w_in / box_h_in
    if ar_i > ar_b: return box_w_in, box_w_in / ar_i
    return box_h_in * ar_i, box_h_in

def add_image_with_callouts(slide, slug, box_x, box_y, box_w, box_h, callouts):
    """Embed a cropped screenshot in the given box, then overlay numbered
    callout circles at fractional (fx, fy) positions on the image.

    callouts: list of (fx, fy) tuples — number = index + 1.
    """
    src = SHOTS / f'{slug}.png'
    if not src.exists():
        # placeholder
        rect(slide, box_x, box_y, box_w, box_h, RULE)
        text(slide, box_x, box_y, box_w, box_h,
             f'[screenshot missing: {slug}]', size=14, color=SLATE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    # crop to box aspect
    crop = SHOTS / f'_um_crop_{slug}.png'
    crop_top(str(src), str(crop), (box_w / box_h))
    iw, ih = fit(str(crop), box_w / 914400, box_h / 914400)
    img_w = Inches(iw); img_h = Inches(ih)
    ix = box_x + (box_w - img_w) // 2
    iy = box_y + (box_h - img_h) // 2
    # frame
    rect(slide, ix - Pt(2), iy - Pt(2), img_w + Pt(4), img_h + Pt(4), RULE)
    slide.shapes.add_picture(str(crop), ix, iy, img_w, img_h)

    # overlays
    dia = Pt(20)
    for n, (fx, fy) in enumerate(callouts, 1):
        cx = ix + int(img_w * fx) - dia // 2
        cy = iy + int(img_h * fy) - dia // 2
        # white halo behind so the circle pops on any background
        halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Pt(2), cy - Pt(2),
                                      dia + Pt(4), dia + Pt(4))
        halo.fill.solid(); halo.fill.fore_color.rgb = WHITE
        halo.line.color.rgb = WHITE; halo.shadow.inherit = False
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, dia, dia)
        circle.fill.solid(); circle.fill.fore_color.rgb = TEAL
        circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
        circle.shadow.inherit = False
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = str(n)
        run.font.name = FONT; run.font.bold = True; run.font.size = Pt(11)
        run.font.color.rgb = WHITE

def what_to_expect(slide, x, y, w, h, title, body):
    round_rect(slide, x, y, w, h, GREEN)
    text(slide, x + Inches(0.18), y + Inches(0.08), w - Inches(0.36), Inches(0.30),
         title.upper(), size=10, bold=True, color=WHITE)
    text(slide, x + Inches(0.18), y + Inches(0.38), w - Inches(0.36), h - Inches(0.5),
         body, size=11, color=WHITE)

# ── Slide skeletons ───────────────────────────────────────────────────────
def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_section_cover(part_no, part_title, kicker_sub):
    s = add_blank()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(3.55), SW, Inches(0.04), TEAL)
    text(s, Inches(0.85), Inches(2.45), Inches(12), Inches(0.45),
         f'PART {part_no}', size=14, bold=True, color=TEAL)
    text(s, Inches(0.85), Inches(2.95), Inches(12), Inches(0.7),
         part_title, size=44, bold=True, color=WHITE)
    text(s, Inches(0.85), Inches(3.75), Inches(12), Inches(2),
         kicker_sub, size=18, color=PAPER)
    notes(s, f'Section divider — Part {part_no}: {part_title}. Pause briefly here when reading aloud.')
    return s

def add_header(slide, kicker, title, route=None):
    rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    text(slide, Inches(0.55), Inches(0.16), Inches(8), Inches(0.30),
         kicker.upper(), size=10, bold=True, color=TEAL)
    text(slide, Inches(0.55), Inches(0.42), Inches(10), Inches(0.5),
         title, size=24, bold=True, color=WHITE)
    if route:
        text(slide, Inches(10.7), Inches(0.46), Inches(2.5), Inches(0.4),
             route, size=10, color=PAPER, align=PP_ALIGN.RIGHT, font='Consolas')

def add_walkthrough_slide(meta):
    """Single-screenshot, numbered-callouts, legend on right, expect strip."""
    s = add_blank()
    add_header(s, meta['kicker'], meta['title'], route=meta.get('route'))

    BX = Inches(0.45); BY = Inches(1.15); BW = Inches(7.55); BH = Inches(5.4)
    add_image_with_callouts(s, meta['slug'], BX, BY, BW, BH, meta['callouts'])

    # right column: numbered legend
    RX = BX + BW + Inches(0.25)
    RY = BY
    RW = SW - RX - Inches(0.4)
    text(s, RX, RY, RW, Inches(0.35),
         'WHAT EACH NUMBER MEANS', size=10, bold=True, color=SLATE)
    numbered(s, RX, RY + Inches(0.35), RW, Inches(4.8),
             meta['legend'], size=11)

    # bottom: what-to-expect strip
    WX = BX; WY = BY + BH + Inches(0.15); WW = SW - Inches(0.85); WH = Inches(0.85)
    what_to_expect(s, WX, WY, WW, WH, 'What to expect',
                   meta.get('expect', 'You should now see the page rendered as above.'))

    notes(s, meta.get('notes', ''))
    return s

def add_diagram_slide(meta):
    """Slides with no embedded screenshot — pure diagram + bullets."""
    s = add_blank()
    add_header(s, meta['kicker'], meta['title'])

    # Left: diagram or hero block. Right: bullets.
    LX = Inches(0.55); LY = Inches(1.15); LW = Inches(6.5); LH = Inches(5.6)
    rect(s, LX, LY, LW, LH, PAPER)

    blocks = meta.get('diagram', [])
    h_each = (LH - Inches(0.4)) / max(len(blocks), 1)
    by = LY + Inches(0.2)
    for label, sub, col in blocks:
        rect(s, LX + Inches(0.3), by, LW - Inches(0.6), h_each - Inches(0.1), col)
        text(s, LX + Inches(0.5), by + Inches(0.12), LW - Inches(1.0), Inches(0.4),
             label, size=15, bold=True, color=WHITE)
        text(s, LX + Inches(0.5), by + Inches(0.5), LW - Inches(1.0), h_each - Inches(0.6),
             sub, size=12, color=WHITE)
        by += h_each

    RX = LX + LW + Inches(0.35); RY = LY
    RW = SW - RX - Inches(0.5)
    text(s, RX, RY, RW, Inches(0.45),
         meta.get('rhead', 'Read me before you click'), size=20, bold=True, color=NAVY)
    bullets(s, RX, RY + Inches(0.65), RW, LH, meta['bullets'], size=13, color=INK)

    notes(s, meta.get('notes', ''))
    return s

# ── Build ─────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ── COVER ────────────────────────────────────────────────────────────────
s = add_blank()
rect(s, 0, 0, SW, SH, NAVY_DARK)
rect(s, 0, Inches(6.1), SW, Inches(0.18), TEAL)
text(s, Inches(0.85), Inches(0.9), Inches(12), Inches(0.4),
     'UGANDA · POINTS OF ENTRY · 2026', size=12, bold=True, color=TEAL)
text(s, Inches(0.85), Inches(1.4), Inches(11.5), Inches(1.6),
     'Web Dashboard', size=58, bold=True, color=WHITE)
text(s, Inches(0.85), Inches(2.7), Inches(11.5), Inches(1.0),
     'User Manual', size=44, color=WHITE)
text(s, Inches(0.85), Inches(4.0), Inches(11.5), Inches(2.0),
     'A complete, click-by-click guide to every screen — written so a first-time user '
     'can sign in, find the right report, drill into the right case, '
     'manage people and places, and recover from any small problem along the way. '
     'Read it cover-to-cover once. Keep it nearby.',
     size=17, color=PAPER)
text(s, Inches(0.85), Inches(6.5), Inches(8), Inches(0.4),
     'Training environment — ug-poe.ecsahc.com/admin · 2026-05-20', size=12, color=PAPER)
text(s, Inches(10.4), Inches(6.5), Inches(2.5), Inches(0.4),
     'Version 1.0', size=12, color=PAPER, align=PP_ALIGN.RIGHT)
notes(s, 'Cover slide. The manual is written so a new user can read it cover-to-cover and operate every part of the dashboard.')

# ── TABLE OF CONTENTS ────────────────────────────────────────────────────
s = add_blank()
rect(s, 0, 0, SW, Inches(0.95), NAVY)
text(s, Inches(0.55), Inches(0.32), Inches(12), Inches(0.55),
     'What is in this manual', size=24, bold=True, color=WHITE)

toc_items = [
    ('Part 1', 'Getting Started', 'Sign in · the sidebar · how every page is built'),
    ('Part 2', 'Quick Reports — your daily reading', 'Eleven one-question reports · drilling into a row · exporting'),
    ('Part 3', 'Acting on Alerts', 'The Alerts hub · acknowledging · the case file deep-dive'),
    ('Part 4', 'Managing People', 'Workforce · the Add Person wizard · per-row actions'),
    ('Part 5', 'Geography', 'Regions · Districts · Hospitals · Countries'),
    ('Part 6', 'PoEs · Annex-1A', 'Registry · Capacity · Status · Roster & Ladder'),
    ('Part 7', 'Reference & troubleshooting', 'Glossary · role card · what to do if'),
]
tx = Inches(0.85); ty = Inches(1.4)
for n, title, sub in toc_items:
    rect(s, tx, ty, Inches(1.4), Inches(0.7), TEAL)
    text(s, tx, ty, Inches(1.4), Inches(0.7), n, size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, tx + Inches(1.65), ty, Inches(10.5), Inches(0.34),
         title, size=17, bold=True, color=NAVY)
    text(s, tx + Inches(1.65), ty + Inches(0.36), Inches(10.5), Inches(0.34),
         sub, size=12, color=SLATE)
    ty += Inches(0.78)
notes(s, 'Table of contents. The manual is sequential — Part 1 first, then read by need.')

# ───────────────────────────────────────────────────────────────────────────
# PART 1 — GETTING STARTED
# ───────────────────────────────────────────────────────────────────────────
add_section_cover(1, 'Getting Started',
    'How to sign in, what you see on first arrival, and the shape every page shares. '
    'Read this part once before you read any other part.')

# 1.1 — How to sign in (annotated login)
add_walkthrough_slide({
    'kicker': 'Part 1 · Getting Started',
    'title':  'Step 1 — Sign in',
    'route':  '/login',
    'slug':   'a01-login',
    'callouts': [
        (0.50, 0.46),   # username field
        (0.50, 0.62),   # password field
        (0.50, 0.78),   # sign-in button
    ],
    'legend': [
        'Type your username here. Your username was given to you on your participant card. '
        'It is not your email — it usually looks like firstname.lastname.',
        'Type your password here. The eye icon on the right shows or hides what you typed — useful if you are not sure.',
        'Click "Sign In" to enter the dashboard. If your username or password is wrong, the page will tell you so in red — '
        'check the spelling and try again.',
    ],
    'expect':
        'After clicking Sign In, the page changes to the Screening Volume report — that is your home screen.',
    'notes':
        'Open with this slide. Tell the trainee: the sign-in page is identical on the training and live dashboards. '
        'The only difference is the URL in the browser bar. Confirm everyone sees the page before continuing.',
})

# 1.2 — Your home screen (Screening Volume landing, annotated)
add_walkthrough_slide({
    'kicker': 'Part 1 · Getting Started',
    'title':  'Step 2 — Your home screen',
    'route':  '/admin/quick-reports/screening-volume',
    'slug':   '00-screening-volume',
    'callouts': [
        (0.04, 0.16),   # sidebar
        (0.50, 0.20),   # page title
        (0.50, 0.40),   # KPI strip
        (0.50, 0.62),   # chart
        (0.95, 0.18),   # top-right user / sign-out area
    ],
    'legend': [
        'The sidebar on the left. Every screen of the dashboard is reached from here.',
        'The page title. Always tells you which screen you are on.',
        'The KPI cards — five to eight headline numbers across the top.',
        'One chart that shows the numbers as a picture. The chart picks its own type for the cleanest view.',
        'Your name, role and a sign-out button — the top-right corner is always you.',
    ],
    'expect':
        'You land here every time you sign in. The numbers refresh on every reload — they are live.',
    'notes':
        'Spend a moment naming each of the five callouts aloud. Make sure trainees can find the sidebar, the page title, '
        'the KPI strip, the chart, and their own name in the top-right.',
})

# 1.3 — The sidebar
add_diagram_slide({
    'kicker': 'Part 1 · Getting Started',
    'title':  'The sidebar — six groups, one rule',
    'diagram': [
        ('Quick Reports',       'Eleven one-question reports. You will use this every day.',          TEAL),
        ('Alert Lifecycle',     'One entry — Alerts. The workflow surface.',                          NAVY),
        ('Workforce',           'One unified page for People · Roles · Assignments.',                 RED),
        ('Geography',           'Reference data — Regions · Districts · Hospitals · Countries.',      AMBER),
        ('PoEs · Annex-1A',     'The borders themselves — Registry · Capacity · Status · Roster.',    GREEN),
        ('Settings',            'Your profile, password, language.',                                  SLATE),
    ],
    'rhead': 'The one rule',
    'bullets': [
        'Your role decides which groups you see. If a group is not in your sidebar, it is not part of your day-to-day work.',
        'Your scope decides which data you see inside each group. National sees everything; PHEOC sees their region; '
        'district sees their district; POE sees their POE.',
        'The dashboard does not say "filtered". The filter is silent — what you see is what you are allowed to see.',
        'Anything hidden is hidden on purpose. It is not broken.',
    ],
    'notes':
        'Most trainees will not see all six groups. That is correct. National Admin sees everything; '
        'a POE-level user sees only Quick Reports and Alerts. Reassure trainees that the menu shape is intentional.',
})

# 1.4 — Anatomy of a page (diagram)
add_diagram_slide({
    'kicker': 'Part 1 · Getting Started',
    'title':  'How every page is built',
    'diagram': [
        ('Filter bar',  'date · POE · status · search — the question you can change.', TEAL),
        ('KPI cards',   '5–8 headline numbers, most important first.',                 NAVY),
        ('One chart',   'A picture of the headline. Reads top-to-bottom.',             AMBER),
        ('One table',   '20 rows on screen. The full list is in the CSV export.',      RED),
    ],
    'rhead': 'Read the page in this order',
    'bullets': [
        'Always read top to bottom. The filter bar tells you what question is being asked. The KPIs tell you the headline numbers. '
        'The chart tells you the picture. The table tells you the names behind the numbers.',
        'Eleven of the screens in this dashboard have this exact shape. Once you know how to read one, you know how to read all eleven.',
        'When you click a row in the table, you usually land in the case file — the deep view for that single record.',
        'Need the full list? The "Export" button at the top right of the table downloads every row to a CSV file.',
    ],
    'notes':
        'This is the most important diagram of the manual. Every Quick Report follows this shape. If a trainee can read one, they can read them all.',
})

# 1.5 — How filters work
add_walkthrough_slide({
    'kicker': 'Part 1 · Getting Started',
    'title':  'How filters work',
    'route':  '(any Quick Report)',
    'slug':   '09-daily-screening',
    'callouts': [
        (0.07, 0.18),   # date range filter
        (0.25, 0.18),   # POE filter
        (0.45, 0.18),   # other filter
        (0.80, 0.18),   # search box
        (0.95, 0.18),   # apply / reset
    ],
    'legend': [
        'Date range — the most-used filter. Defaults to the last 7 days.',
        'POE selector — pick one POE to focus on it. Leave blank to see all your POEs.',
        'Status / risk / category filters — the page shows whichever of these are relevant.',
        'Free-text search — typing here narrows the table to matching rows.',
        'Apply or Reset. Filters update both the KPIs and the chart, not only the table.',
    ],
    'expect':
        'Click Apply and the whole page redraws. The URL in your browser updates too — copy the URL to share a filtered view.',
    'notes':
        'Stress: filters update everything, not only the table. KPIs, chart and table all redraw together.',
})

# 1.6 — Reading the UI language (badges / pills / colors)
add_diagram_slide({
    'kicker': 'Part 1 · Getting Started',
    'title':  'The colour language',
    'diagram': [
        ('CRITICAL · HIGH',  'Red and dark red — needs eyes now.',          RED),
        ('MEDIUM',           'Amber — keep an eye on it.',                  AMBER),
        ('LOW · INFO',       'Teal and grey — routine.',                    TEAL),
        ('CLOSED · DONE',    'Green — the case has been resolved.',         GREEN),
    ],
    'rhead': 'Pills and badges',
    'bullets': [
        'A "pill" is a small coloured rounded label. Status, risk and tier all appear as pills.',
        'Red and dark red always mean urgent. Amber always means caution. Teal and grey always mean routine.',
        'Green always means a case has been closed or an action has been completed.',
        'The same colour means the same thing on every screen — there are no exceptions.',
    ],
    'notes': 'Teach the colour language once. It applies everywhere.',
})

# ───────────────────────────────────────────────────────────────────────────
# PART 2 — QUICK REPORTS
# ───────────────────────────────────────────────────────────────────────────
add_section_cover(2, 'Quick Reports — your daily reading',
    'Eleven reports that each answer one question. They all share the same shape — '
    'so once you can read one, you can read all eleven.')

QR_LEGEND_TEMPLATE = [
    'Filter bar — change the date range, the POE and other filters here.',
    'KPI cards — headline numbers, left-to-right by importance.',
    'The chart — a picture of the headline. Read the axis labels first.',
    'The table — 20 rows on screen, full list in the Export button.',
]

QR_CALLOUTS = [(0.50, 0.18), (0.50, 0.34), (0.50, 0.60), (0.50, 0.88)]

QR_SURFACES = [
    {
        'slug':  '00-screening-volume',
        'title': 'Screening Volume',
        'route': '/admin/quick-reports/screening-volume',
        'one_line': 'How many travellers, split by primary, secondary, gender and age.',
        'who':     'Everyone. This is your home screen.',
        'expect':  'Lands on Past 7 days by default. Change the date range, the chart and KPIs update live.',
    },
    {
        'slug':  '01-suspected-cases',
        'title': 'Suspected Cases',
        'route': '/admin/quick-reports/suspected-cases',
        'one_line': 'Who do we suspect right now, of what disease, at what risk?',
        'who':     'Supervisors and the surveillance team — read this first thing every morning.',
        'expect':  'Click any row to open the case file for that traveller. Use this as your daily working list.',
    },
    {
        'slug':  '02-confirmed-cases',
        'title': 'Confirmed Cases',
        'route': '/admin/quick-reports/confirmed-cases',
        'one_line': 'Of the suspects, who has been confirmed, probable, or ruled out by the lab?',
        'who':     'Surveillance and lab liaisons.',
        'expect':  'Read this together with Suspected Cases — one is what we think, the other is what the lab says.',
    },
    {
        'slug':  '03-alert-database',
        'title': 'Alert Database',
        'route': '/admin/quick-reports/alert-database',
        'one_line': 'Every alert ever — the full ledger, with owners.',
        'who':     'Anyone auditing alert history. Use this when you need the full list.',
        'expect':  'Keep the date filter to 7 or 30 days while browsing — All-time loads thousands of rows.',
    },
    {
        'slug':  '04-alert-analysis',
        'title': 'Alert Analysis',
        'route': '/admin/quick-reports/alert-analysis',
        'one_line': 'Where the HIGH-risk and tier-1 alerts are clustering.',
        'who':     'PHEOC and supervisors — open at the weekly briefing.',
        'expect':  'The on-screen table shows high-stakes rows only. The full cohort is in the CSV export.',
    },
    {
        'slug':  '05-alert-outcomes',
        'title': 'Alert Outcomes',
        'route': '/admin/quick-reports/alert-outcomes',
        'one_line': 'How quickly we acknowledge and close alerts. SLA performance.',
        'who':     'Supervisors and leadership.',
        'expect':  'Two numbers matter: median acknowledge minutes and median close hours. Read them first.',
    },
    {
        'slug':  '06-symptom-spread',
        'title': 'Symptom Spread',
        'route': '/admin/quick-reports/symptom-spread',
        'one_line': 'Which symptoms are appearing right now, and which are red flags.',
        'who':     'Surveillance and clinical leads.',
        'expect':  'Red-flag symptoms appear in red automatically. You do not need to configure anything.',
    },
    {
        'slug':  '07-poe-analysis',
        'title': 'POE Analysis',
        'route': '/admin/quick-reports/poe-analysis',
        'one_line': 'Which borders are busy, which are dark, which produce the most alerts.',
        'who':     'Anyone comparing POEs side-by-side.',
        'expect':  '"Dark" means no recent activity. That is information, not necessarily a problem.',
    },
    {
        'slug':  '08-country-analysis',
        'title': 'Country Analysis',
        'route': '/admin/quick-reports/country-analysis',
        'one_line': 'Where travellers are from, where they have been, and where they have transited.',
        'who':     'Surveillance — travel-history questions.',
        'expect':  'Endemic-country flow is highlighted in red automatically.',
    },
    {
        'slug':  '09-daily-screening',
        'title': 'Daily Screening',
        'route': '/admin/quick-reports/daily-screening',
        'one_line': 'How many travellers per day, and what fraction became referrals.',
        'who':     'Daily operations standup.',
        'expect':  'Compare today vs yesterday vs the 7-day average — those three KPIs are at the top.',
    },
    {
        'slug':  '10-user-analysis',
        'title': 'User Analysis',
        'route': '/admin/quick-reports/user-analysis',
        'one_line': 'Which officers are active, which are dormant, and who is carrying the load.',
        'who':     'Monthly workforce reviews.',
        'expect':  'The top performer is highlighted in red — the eye finds them first.',
    },
]

# Intro slide before the eleven
add_diagram_slide({
    'kicker': 'Part 2 · Quick Reports',
    'title':  'Eleven reports, one shape',
    'diagram': [
        ('1. Filter bar',  'date · POE · status · search — change the question.',  TEAL),
        ('2. KPI cards',   '5–8 headline numbers across the top.',                  NAVY),
        ('3. The chart',   'A picture of the headline. The chart picks its type.', AMBER),
        ('4. The table',   '20 rows on screen, full list in Export.',               RED),
    ],
    'rhead': 'What follows',
    'bullets': [
        'The next eleven slides each open one Quick Report.',
        'Each slide shows the page, what each part is, and what to expect when you arrive.',
        'You can come back here any time as a refresher.',
    ],
    'notes': 'Intro slide before the eleven Quick Reports. Anchors the trainee in the shared shape before walking each one.',
})

for qr in QR_SURFACES:
    add_walkthrough_slide({
        'kicker': 'Part 2 · Quick Reports',
        'title':  qr['title'],
        'route':  qr['route'],
        'slug':   qr['slug'],
        'callouts': QR_CALLOUTS,
        'legend':   QR_LEGEND_TEMPLATE,
        'expect':   qr['expect'],
        'notes':
            f'{qr["title"]}. The one question this answers: {qr["one_line"]} '
            f'Who opens it: {qr["who"]}',
    })

# Drilling into a row → the case file
add_walkthrough_slide({
    'kicker': 'Part 2 · Quick Reports',
    'title':  'Drilling into a row — the case file',
    'route':  '/admin/alerts/{id}/case-file',
    'slug':   'a02-case-file',
    'callouts': [
        (0.50, 0.06),   # title
        (0.50, 0.20),   # tabs
        (0.50, 0.45),   # main content panel
        (0.85, 0.10),   # status badge
    ],
    'legend': [
        'The page title shows the alert code and the traveller name.',
        'Six tabs across the top — Patient · Travel · Clinical · Timeline · Communications · Outcome.',
        'The main panel shows every field captured on the phone. Switch tabs to switch sections.',
        'Status badge at the top-right tells you whether the case is OPEN, ACKNOWLEDGED, CLOSED or REOPENED.',
    ],
    'expect':
        'This is the deepest view in the dashboard — 100 % of what was captured on the phone is here.',
    'notes':
        'The case file is opened from a row in any Quick Report and from the Alerts hub. '
        'Click anywhere in the row, or click the 📋 icon on the right of the row.',
})

# Export to CSV
add_diagram_slide({
    'kicker': 'Part 2 · Quick Reports',
    'title':  'Exporting to CSV',
    'diagram': [
        ('Step 1', 'Apply your filters so the table shows what you want.', TEAL),
        ('Step 2', 'Click the "Export CSV" button at the top right of the table.', NAVY),
        ('Step 3', 'Your browser downloads a .csv file with every matching row — not only the 20 on screen.', AMBER),
        ('Step 4', 'Open the .csv in Excel, LibreOffice, or Google Sheets.', GREEN),
    ],
    'rhead': 'About the export',
    'bullets': [
        'The export carries the full cohort — every row that matches your filters, not only the 20 rows on screen.',
        'The columns in the CSV match the columns in the on-screen table.',
        'Time stamps in the CSV use Africa/Kampala (UTC+3) — the same as the screen.',
        'Large exports can take a few seconds. Keep the page open until the download finishes.',
    ],
    'notes': 'The export is the answer to "where are the other rows?" — keep this slide as a reference.',
})

# ───────────────────────────────────────────────────────────────────────────
# PART 3 — ACTING ON ALERTS
# ───────────────────────────────────────────────────────────────────────────
add_section_cover(3, 'Acting on Alerts',
    'How to read the alert list, how to acknowledge, how to open the case file. '
    'Supervisors and PHEOC act here — everyone else reads.')

# 3.1 — Alerts hub annotated
add_walkthrough_slide({
    'kicker': 'Part 3 · Acting on Alerts',
    'title':  'The Alerts hub',
    'route':  '/admin/alerts',
    'slug':   '11-alerts-hub',
    'callouts': [
        (0.30, 0.15),  # tab New
        (0.50, 0.15),  # tabs row
        (0.50, 0.32),  # filter bar
        (0.50, 0.55),  # table
        (0.93, 0.55),  # action icons on a row
    ],
    'legend': [
        'The "New" tab — these are the alerts that need acknowledgement. This is the default landing.',
        'Tabs across the top — New · Being worked on · Closed · Reopened · All. Click a tab to switch the view.',
        'Filters underneath the tabs — date window, risk, response team, district, POE, search.',
        'The table — each row is one alert. The columns show traveller, risk pill, status pill, owner, when raised.',
        'Action icons on the right — click ✓ to acknowledge, click 📋 to open the case file.',
    ],
    'expect':
        'The "New" tab number on the chip tells you how many alerts are waiting. Open the case file before acknowledging.',
    'notes':
        'POE-level users cannot acknowledge — the ✓ icon is hidden for them. They can only open the case file. '
        'District Supervisors, PHEOC and National can act on alerts.',
})

# 3.2 — Acknowledging an alert (step-by-step)
add_diagram_slide({
    'kicker': 'Part 3 · Acting on Alerts',
    'title':  'How to acknowledge an alert',
    'diagram': [
        ('Step 1', 'On the Alerts hub, find the alert in the "New" tab.', TEAL),
        ('Step 2', 'Click the 📋 icon in the row\'s action strip — the case file opens.', NAVY),
        ('Step 3', 'Read the case file. Check Patient · Travel · Clinical tabs.', AMBER),
        ('Step 4', 'Return to the Alerts hub. Click ✓ on the same row.', GREEN),
    ],
    'rhead': 'After acknowledging',
    'bullets': [
        'The alert moves from the "New" tab to the "Being worked on" tab.',
        'Your name is recorded as the owner. The acknowledged timestamp is recorded too.',
        'The alert stays in "Being worked on" until you (or another supervisor) close it.',
        'You can re-open a closed alert from the "Closed" tab if new information arrives.',
    ],
    'notes':
        'Stress: always open the case file FIRST. Acknowledging without reading the case is a recordable mistake.',
})

# 3.3 — The case file deep view
add_walkthrough_slide({
    'kicker': 'Part 3 · Acting on Alerts',
    'title':  'The case file — 100% of what was captured',
    'route':  '/admin/alerts/{id}/case-file',
    'slug':   'a02-case-file',
    'callouts': [
        (0.10, 0.10),   # back link
        (0.50, 0.20),   # tabs row
        (0.10, 0.45),   # patient identity
        (0.50, 0.65),   # clinical / actions panel
        (0.93, 0.10),   # status badge top right
    ],
    'legend': [
        'The "Back" link returns you to the previous list (Quick Report or Alerts hub).',
        'Six tabs — Patient · Travel · Clinical · Timeline · Communications · Outcome.',
        'Patient identity panel — name, age, gender, nationality, document.',
        'The active tab\'s content. Switch tabs to see different aspects of the case.',
        'The status badge — OPEN, ACKNOWLEDGED, CLOSED or REOPENED at a glance.',
    ],
    'expect':
        'Every field you can see here was captured on the screener\'s phone. Nothing is added in the dashboard.',
    'notes':
        'The case file is the canonical "100% of mobile data" view. If a field exists in the mobile app, it shows here.',
})

# ───────────────────────────────────────────────────────────────────────────
# PART 4 — MANAGING PEOPLE (Workforce)
# ───────────────────────────────────────────────────────────────────────────
add_section_cover(4, 'Managing People',
    'The Workforce page is where users are created, given roles, and assigned to POEs. '
    'Only National Admin can create new users — everyone else reads their scope.')

# 4.1 — Workforce overview
add_walkthrough_slide({
    'kicker': 'Part 4 · Managing People',
    'title':  'Workforce — one unified page',
    'route':  '/admin/workforce',
    'slug':   '12-workforce',
    'callouts': [
        (0.30, 0.16),   # tabs row
        (0.93, 0.16),   # + Add person button
        (0.50, 0.32),   # filter bar
        (0.50, 0.60),   # table
        (0.90, 0.60),   # per-row icons
    ],
    'legend': [
        'Three tabs across the top — Users · Roles · Assignments. Click to switch.',
        '+ Add person button — opens the four-step wizard. National Admin only.',
        'Filter bar — role, POE, status, search. Combine filters to narrow the list.',
        'The users table — name, username, role, jurisdiction, status, last sign-in.',
        'Per-row action icons — edit, suspend, reset password, regenerate invite, disable.',
    ],
    'expect':
        'This page shows 61 users on the training server. The "Users" tab is the default — Roles and Assignments are in the other two tabs.',
    'notes':
        'If you are not National Admin, the + Add person button may be greyed out. Read-only walks are still useful.',
})

# 4.2 — The Add Person wizard
add_walkthrough_slide({
    'kicker': 'Part 4 · Managing People',
    'title':  'The Add Person wizard',
    'route':  '/admin/workforce → + Add person',
    'slug':   'a04-workforce-wizard',
    'callouts': [
        (0.50, 0.10),   # wizard title
        (0.50, 0.20),   # step progress dots
        (0.50, 0.42),   # current step content
        (0.30, 0.85),   # back button
        (0.70, 0.85),   # next button
    ],
    'legend': [
        'The wizard title — confirms what you are about to do.',
        'Step progress dots — four steps in total. The current step is highlighted.',
        'The form for the current step. Fill in every required field before you can advance.',
        'Back — go to the previous step without losing what you have entered.',
        'Next — advance to the following step. Becomes "Create person" on the final step.',
    ],
    'expect':
        'The wizard is atomic. Nothing is saved until you click "Create person" on Step 4. Closing the wizard discards everything.',
    'notes':
        'Always tell the trainee: nothing is saved until the final Create person click. They can safely close the wizard.',
})

# 4.3, 4.4, 4.5, 4.6 — Wizard steps
for step in [
    {
        'title': 'Wizard step 1 — Identity',
        'rhead': 'What this step asks',
        'bullets': [
            'Full name — the person\'s real name as they would sign a letter.',
            'Username — what they will type into the sign-in box. Lower-case, no spaces, usually firstname.lastname.',
            'Email — must be unique across the whole country, not only your district.',
            'Phone (optional) — used for SMS notifications.',
            'When you click Next, the system checks that the username and email are not already taken.',
        ],
        'diagram': [
            ('Full name', 'Real name. "Sarah Nakimuli".',                    TEAL),
            ('Username',  'firstname.lastname, no spaces.',                  NAVY),
            ('Email',     'Unique across the whole country.',                AMBER),
            ('Phone',     'Optional — for SMS notifications.',               SLATE),
        ],
    },
    {
        'title': 'Wizard step 2 — Role',
        'rhead': 'Which role to pick',
        'bullets': [
            'The role decides what the new user sees in the dashboard and on the mobile app.',
            'NATIONAL_ADMIN — super-user. Sees everything, can do everything.',
            'PHEOC_OFFICER — regional coordination. Sees their region\'s districts and POEs.',
            'DISTRICT_SUPERVISOR — sees their district\'s POEs.',
            'POE_PRIMARY / POE_SECONDARY — screeners. See only their POE.',
            'Click the radio card next to the role you want, then Next.',
        ],
        'diagram': [
            ('NATIONAL_ADMIN',     'Sees and does everything.',                                  RED),
            ('PHEOC_OFFICER',      'Regional coordination — multi-district scope.',              AMBER),
            ('DISTRICT_SUPERVISOR','District scope — all POEs in the district.',                 NAVY),
            ('POE_PRIMARY / SECONDARY', 'Single-POE screeners.',                                 TEAL),
        ],
    },
    {
        'title': 'Wizard step 3 — Jurisdiction',
        'rhead': 'Where the person works',
        'bullets': [
            'This step changes depending on the role you picked in Step 2.',
            'For NATIONAL — no jurisdiction is asked, the country is implied.',
            'For PHEOC — pick a region. The system will show only the regions you are allowed to assign.',
            'For DISTRICT — pick a district. The region is inferred from the district.',
            'For POE — pick a POE. The district and region are inferred from the POE.',
        ],
        'diagram': [
            ('National',  'No jurisdiction asked.',                                              SLATE),
            ('Region',    'Pick one region — for PHEOC officers.',                               NAVY),
            ('District',  'Pick one district — for District Supervisors.',                       AMBER),
            ('POE',       'Pick one POE — for POE-level screeners.',                             TEAL),
        ],
    },
    {
        'title': 'Wizard step 4 — Invite mode',
        'rhead': 'How the new user signs in for the first time',
        'bullets': [
            'Two ways to invite a new user. Pick one.',
            'Credential — the wizard generates a temporary password and shows it to you on screen. '
            'Write it down or copy it before you close the dialog. The user must change it on first sign-in.',
            'Email invite — the wizard sends an email with a one-time link. The link expires in 7 days. '
            'The user creates their own password when they click the link.',
            'Click "Create person" to finish. The new user appears at the top of the Users tab.',
        ],
        'diagram': [
            ('Credential',  'Generate a password, hand it over directly.',  TEAL),
            ('Email invite','Send a 7-day one-time link.',                   NAVY),
        ],
    },
]:
    add_diagram_slide({
        'kicker': 'Part 4 · Managing People',
        'title':  step['title'],
        'diagram': step['diagram'],
        'rhead':   step['rhead'],
        'bullets': step['bullets'],
        'notes':   f'{step["title"]} walkthrough. Read the bullets aloud, point at each block in the diagram.',
    })

# 4.7 — Per-row icons
add_diagram_slide({
    'kicker': 'Part 4 · Managing People',
    'title':  'Per-user actions — the icons on each row',
    'diagram': [
        ('Edit',          'Change the user\'s name, email, phone, role.',                  TEAL),
        ('Suspend',       'Block sign-in temporarily without deleting the user.',          AMBER),
        ('Reset password','Issue a new temporary password. User must change on next sign-in.', NAVY),
        ('Disable',       'Soft-delete. The user is removed from active lists but history is preserved.', RED),
    ],
    'rhead': 'When to use each',
    'bullets': [
        'Edit — used when a person\'s name changes, or they move to a new POE.',
        'Suspend — used when a person is on leave. Easily reversed.',
        'Reset password — used when the user has forgotten their password.',
        'Disable — used when a person has left the organisation. Their history stays in reports.',
    ],
    'notes': 'These four icons cover 95 % of day-to-day user management.',
})

# ───────────────────────────────────────────────────────────────────────────
# PART 5 — GEOGRAPHY
# ───────────────────────────────────────────────────────────────────────────
add_section_cover(5, 'Geography — the reference hierarchy',
    'Country · Region · District · Hospital. The reference data behind everything else. '
    'Only National Admin writes here; everyone else reads.')

GEO = [
    {
        'slug':  '13-geo-regions',
        'title': 'Regions',
        'route': '/admin/geo/provinces',
        'one_line': 'The country\'s regional PHEOCs. The first layer under the country.',
    },
    {
        'slug':  '14-geo-districts',
        'title': 'Districts',
        'route': '/admin/geo/districts',
        'one_line': 'Border-adjacent districts. Each district sits under one region.',
    },
    {
        'slug':  '15-geo-hospitals',
        'title': 'Hospitals',
        'route': '/admin/geo/hospitals',
        'one_line': 'The referral hospital roster. Each hospital sits under one district.',
    },
    {
        'slug':  '16-geo-countries',
        'title': 'Countries',
        'route': '/admin/geo/countries',
        'one_line': 'The country itself — Uganda. A single row.',
    },
]
for g in GEO:
    add_walkthrough_slide({
        'kicker': 'Part 5 · Geography',
        'title':  g['title'],
        'route':  g['route'],
        'slug':   g['slug'],
        'callouts': [
            (0.50, 0.18),   # filters
            (0.50, 0.34),   # toolbar (search + +new)
            (0.50, 0.60),   # table
            (0.93, 0.60),   # per-row actions
        ],
        'legend': [
            'Filters — country, parent geo, status, search.',
            'Toolbar — "+ New" button on the right, status tabs on the left.',
            'The table — code, name, status, count of children (e.g. districts in this region).',
            'Per-row actions — edit, soft-delete, restore. Visible to National Admin only.',
        ],
        'expect': g['one_line'],
        'notes':
            'For Regions and Districts: renaming a row cascades to every POE underneath. '
            'For Hospitals: greenfield, the table is intentionally light today. '
            'For Countries: single row, only the metadata is editable.',
    })

# Cascading rename warning
add_diagram_slide({
    'kicker': 'Part 5 · Geography',
    'title':  'The cascading rename rule',
    'diagram': [
        ('Region renames',   'Every POE in that region updates its "region" field.',  AMBER),
        ('District renames', 'Every POE in that district updates its "district" field.', AMBER),
        ('Bundle bumps',     'A geo bundle-version is bumped so phones re-sync.',     NAVY),
        ('Cannot delete',    'A region or district with PoEs/hospitals under it cannot be deleted.', RED),
    ],
    'rhead': 'What this means for you',
    'bullets': [
        'Renaming is a real operational event — every screener device picks up the new name on the next sync.',
        'Avoid renaming during the working day if possible. Pick a low-traffic window.',
        'If you must delete a region or district, move its children first.',
        'When in doubt, ask National Admin before clicking save.',
    ],
    'notes': 'This is the only "be careful" rule in Geography. Repeat it.',
})

# ───────────────────────────────────────────────────────────────────────────
# PART 6 — PoEs · Annex-1A
# ───────────────────────────────────────────────────────────────────────────
add_section_cover(6, 'PoEs · Annex-1A — the borders',
    'Four operational surfaces for each Point of Entry. Who they are. How capable. Open or closed. Who to call.')

# PoE Registry
add_walkthrough_slide({
    'kicker': 'Part 6 · PoEs · Annex-1A',
    'title':  'PoE Registry',
    'route':  '/admin/geo/poes',
    'slug':   '17-poe-registry',
    'callouts': [
        (0.50, 0.18),   # filter bar
        (0.93, 0.18),   # + New PoE
        (0.50, 0.55),   # table
        (0.93, 0.55),   # per-row icons
    ],
    'legend': [
        'Filter bar — country, region, district, PoE type, transport mode, status, search.',
        '+ New PoE button — opens the five-step wizard. National Admin only.',
        'The table — every PoE, with code, name, type, region, district, status.',
        'Per-row icons — edit, soft-delete, restore.',
    ],
    'expect':
        'There are 39 PoEs on the training server. The table is paginated — 25 rows per page.',
    'notes':
        'PoEs are the operational anchor of the system. Every alert eventually belongs to a PoE.',
})

# PoE Registry wizard
add_walkthrough_slide({
    'kicker': 'Part 6 · PoEs · Annex-1A',
    'title':  'PoE Registry — the New PoE wizard',
    'route':  '/admin/geo/poes → + New PoE',
    'slug':   'a03-poe-wizard',
    'callouts': [
        (0.50, 0.08),   # wizard title
        (0.50, 0.18),   # step dots
        (0.50, 0.40),   # current step form
        (0.30, 0.85),   # back
        (0.70, 0.85),   # next
    ],
    'legend': [
        'Wizard title — confirms you are creating a new PoE.',
        'Five step dots — Identity · Location · Flags · Context · Review.',
        'The current step form. Type the PoE name; the system suggests its type and transport mode.',
        'Back — return to the previous step without losing data.',
        'Next — advance. Becomes "Create" on the Review step.',
    ],
    'expect':
        'When you start typing a PoE name, suggestion chips appear underneath — pick one and several fields auto-fill.',
    'notes':
        'The five-step wizard takes about a minute to fill in. The "Review" step shows you everything before save.',
})

# PoE Capacity
add_walkthrough_slide({
    'kicker': 'Part 6 · PoEs · Annex-1A',
    'title':  'PoE Capacity — Annex-1A scoring',
    'route':  '/admin/poe/capacity',
    'slug':   '18-poe-capacity',
    'callouts': [
        (0.20, 0.18),   # DRAFT tab
        (0.50, 0.18),   # SUBMITTED tab
        (0.80, 0.18),   # REVIEWED tab
        (0.50, 0.55),   # table
    ],
    'legend': [
        'DRAFT tab — assessments still being filled in. Editable.',
        'SUBMITTED tab — assessments handed in for review. Read-only.',
        'REVIEWED tab — finalised assessments. Locked.',
        'The table — one row per assessment, with the overall 0–100 score.',
    ],
    'expect':
        'The workflow is one-way: DRAFT → SUBMITTED → REVIEWED. Once you submit, you cannot edit unless a reviewer kicks it back.',
    'notes':
        'Eight scoring dimensions per assessment. Each dimension is 1–5. The overall score is computed automatically.',
})

# PoE Status
add_walkthrough_slide({
    'kicker': 'Part 6 · PoEs · Annex-1A',
    'title':  'PoE Status',
    'route':  '/admin/poe/status',
    'slug':   '19-poe-status',
    'callouts': [
        (0.50, 0.18),   # filter
        (0.50, 0.40),   # current status per POE
        (0.50, 0.70),   # recent log
        (0.93, 0.40),   # action
    ],
    'legend': [
        'Filter by POE, by district, by status.',
        'Current status — the latest row per POE. Shown as a large pill.',
        'Recent log — the last 30 status changes across all POEs in your scope.',
        'Record-next button — change a POE\'s status (you cannot "edit" a current status, you record the next one).',
    ],
    'expect':
        'Five possible statuses: OPEN · CLOSED · REDUCED_HOURS · EMERGENCY_CLOSED · MAINTENANCE.',
    'notes':
        'Posting a new status automatically closes the previous one. You record next, you do not edit current.',
})

# Roster & Ladder
add_walkthrough_slide({
    'kicker': 'Part 6 · PoEs · Annex-1A',
    'title':  'Roster & Ladder — who to call',
    'route':  '/admin/poe/contacts',
    'slug':   '20-roster-ladder',
    'callouts': [
        (0.50, 0.18),   # filters
        (0.50, 0.42),   # table — contact rows
        (0.85, 0.42),   # receives flags
        (0.93, 0.42),   # escalates-to chain
    ],
    'legend': [
        'Filter chips — level (POE / District / PHEOC / National / WHO), preferred channel, search.',
        'The contact list — name, level, organisation, phone, email.',
        'Receives flags — ten yes/no flags showing what kinds of alerts this contact gets.',
        'Escalates-to — who is next on the ladder if this contact does not respond.',
    ],
    'expect':
        'When an alert is raised, it walks the ladder from POE upward until someone responds.',
    'notes':
        'Ten flags per contact. Five levels in the ladder. Maximum chain depth is 5 hops.',
})

# ───────────────────────────────────────────────────────────────────────────
# PART 7 — REFERENCE & TROUBLESHOOTING
# ───────────────────────────────────────────────────────────────────────────
add_section_cover(7, 'Reference & troubleshooting',
    'The pages you keep nearby. Glossary, role card, what-to-do-if.')

# Troubleshooting
add_diagram_slide({
    'kicker': 'Part 7 · Reference',
    'title':  'What to do if…',
    'diagram': [
        ('Page is blank',           'Refresh the page. If still blank, sign out and sign back in.',     AMBER),
        ('Filter does not stick',   'Click Apply again. If it still resets, refresh the page.',         AMBER),
        ('Export is slow',          'Wait. Large exports take a few seconds. Do not close the page.',   AMBER),
        ('You cannot see X',        'Your role decides what you see. If X is hidden, it is not yours.', NAVY),
        ('Error in red',            'Read the message. Most errors tell you exactly what to fix.',      RED),
        ('Lost your password',      'Ask your admin to reset it from the Workforce page.',              SLATE),
    ],
    'rhead': 'Three rules to remember',
    'bullets': [
        'Refresh first. Many small UI hiccups disappear on a refresh.',
        'Read the message. The dashboard is verbose on errors — it usually tells you the cause.',
        'When in doubt, sign out and sign back in. A fresh session fixes most session-state issues.',
    ],
    'notes': 'Three rules cover most user-reported problems.',
})

# Glossary
s = add_blank()
add_header(s, 'Part 7 · Reference', 'Glossary')
GLOSS = [
    ('Acknowledge',  'Mark an alert as "being worked on". Done by supervisors and above.'),
    ('Alert',        'A record raised when a secondary case is closed as SUSPECTED or CONFIRMED.'),
    ('Bundle',       'The reference data set the mobile app downloads. Bumps when geo data changes.'),
    ('Case file',    'The deep view of one case — every field captured on the phone.'),
    ('CRUD',         'Create · Read · Update · Delete. The four basic data operations.'),
    ('CSV',          'A spreadsheet-friendly text file. The format used for all exports.'),
    ('IHR tier',     'WHO classification — Tier 1 (always notifiable), Tier 2 (Annex-2), Tier 3 (routine).'),
    ('KPI',          'Key Performance Indicator. The headline number on a card.'),
    ('PHEOC',        'Public Health Emergency Operations Centre. A region\'s coordination unit.'),
    ('Pill',         'A small coloured rounded label. Status, risk and tier all use pills.'),
    ('PoE',          'Point of Entry — airport, port, land border, rail crossing.'),
    ('Quick Report', 'A one-question report. Eleven exist; all share the same shape.'),
    ('RBAC',         'Role-Based Access Control. Your role decides your menu.'),
    ('Referral',     'A primary screening that became a secondary case because of symptoms.'),
    ('Scope',        'The geographic area you are allowed to see — country, region, district, or POE.'),
    ('SLA',          'Service Level Agreement. Targets for acknowledge and close times.'),
    ('Soft-delete',  'A "deletion" that preserves history but hides the row from active lists.'),
    ('Suspected',    'A case the engine flagged as possibly a specific disease.'),
    ('Tab',          'A horizontal switch at the top of a page that changes what the page shows.'),
    ('Tier',         'See "IHR tier".'),
]
# two columns
col_w = (SW - Inches(1.5)) / 2
half = (len(GLOSS) + 1) // 2
left  = GLOSS[:half]; right = GLOSS[half:]
def put_col(x, items):
    tb = s.shapes.add_textbox(x, Inches(1.2), col_w, Inches(6))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, (term, defn) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.30
        rT = p.add_run(); rT.text = f'{term}  '
        rT.font.name = FONT; rT.font.size = Pt(12); rT.font.bold = True; rT.font.color.rgb = NAVY
        rD = p.add_run(); rD.text = defn
        rD.font.name = FONT; rD.font.size = Pt(11); rD.font.color.rgb = INK
put_col(Inches(0.55), left)
put_col(Inches(0.55) + col_w + Inches(0.4), right)
notes(s, 'Glossary — keep this page bookmarked. Twenty terms cover everything in this manual.')

# Role reference card
add_diagram_slide({
    'kicker': 'Part 7 · Reference',
    'title':  'Role reference — who sees what',
    'diagram': [
        ('NATIONAL_ADMIN',      'Everything in the dashboard. Only role that can create users.',        RED),
        ('PHEOC_OFFICER',       'Their region\'s districts and POEs. Acts on alerts in scope.',         AMBER),
        ('DISTRICT_SUPERVISOR', 'Their district\'s POEs. Acts on alerts in scope.',                     NAVY),
        ('POE_PRIMARY',         'Their POE only. Reads everything in scope, captures screenings on the phone.', TEAL),
        ('POE_SECONDARY',       'Their POE only. Walks secondary cases on the phone.',                  TEAL),
        ('POE_DATA_OFFICER',    'Their POE only. Reads + submits routine reports.',                     SLATE),
    ],
    'rhead': 'How to read your role',
    'bullets': [
        'Your role badge is on the top-right of every page next to your name.',
        'If a menu item is missing, your role does not need it. The dashboard hides things on purpose.',
        'If you change roles within the organisation, ask National Admin to update your role from the Workforce page.',
    ],
    'notes': 'Six roles total. Keep this slide as a reference.',
})

# Closing slide
s = add_blank()
rect(s, 0, 0, SW, SH, NAVY_DARK)
rect(s, 0, Inches(6.1), SW, Inches(0.18), TEAL)
text(s, Inches(0.85), Inches(2.6), Inches(12), Inches(0.5),
     'END OF MANUAL', size=14, bold=True, color=TEAL)
text(s, Inches(0.85), Inches(3.1), Inches(12), Inches(1.6),
     'You are ready.', size=52, bold=True, color=WHITE)
text(s, Inches(0.85), Inches(4.6), Inches(12), Inches(1.2),
     'Sign in. Open Screening Volume. Read the headline numbers. '
     'Open one row, walk the case file. That is one full circuit.',
     size=18, color=PAPER)
text(s, Inches(0.85), Inches(6.5), Inches(8), Inches(0.4),
     'Keep this manual nearby. Refer to it whenever a screen is unfamiliar.',
     size=12, color=PAPER)
text(s, Inches(10.4), Inches(6.5), Inches(2.5), Inches(0.4),
     '— end —', size=12, color=PAPER, align=PP_ALIGN.RIGHT)
notes(s, 'Closing slide. Read the call-to-action.')

# ── Save ─────────────────────────────────────────────────────────────────
OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPTX)
print(f'Wrote {OUT_PPTX}')
print(f'Slides: {len(prs.slides)}')
print(f'Size:   {OUT_PPTX.stat().st_size / 1024:.0f} KB')
